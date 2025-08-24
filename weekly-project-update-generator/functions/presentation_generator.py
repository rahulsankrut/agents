"""
Presentation Generator using python-pptx

This module handles the creation of PowerPoint presentations from analyzed
images and generated content.
"""

import logging
from typing import Dict, Any, List
from io import BytesIO
import base64
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """
    Generates PowerPoint presentations from project data and image analyses
    """
    
    def __init__(self):
        """Initialize the PresentationGenerator"""
        self.slide_width = Inches(13.33)  # Standard 16:9 aspect ratio
        self.slide_height = Inches(7.5)
        self.margin = Inches(0.5)
        self.content_width = self.slide_width - (2 * self.margin)
        self.content_height = self.slide_height - (2 * self.margin)
    
    def generate_presentation(self, project_details: Dict[str, Any], image_analyses: List[Dict[str, Any]]) -> bytes:
        """
        Generate a complete PowerPoint presentation
        
        Args:
            project_details: Project information
            image_analyses: List of analyzed images and content
            
        Returns:
            PowerPoint file as bytes
        """
        try:
            # Create new presentation
            prs = Presentation()
            
            # Set slide dimensions
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height
            
            # Generate slides
            self._create_title_slide(prs, project_details)
            self._create_summary_slide(prs, project_details, image_analyses)
            
            # Create content slides for each image
            for analysis in image_analyses:
                self._create_content_slide(prs, analysis)
            
            # Create conclusion slide
            self._create_conclusion_slide(prs, project_details)
            
            # Save to bytes
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            
            logger.info(f"Presentation generated successfully with {len(prs.slides)} slides")
            return output.getvalue()
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise
    
    def _create_title_slide(self, prs: Presentation, project_details: Dict[str, Any]):
        """Create the title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background shape
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.slide_width, self.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(44, 62, 80)  # Dark blue-gray
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin, Inches(2),
            self.content_width, Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Weekly Project Update"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(48)
        title_run.font.color.rgb = RGBColor(255, 255, 255)
        title_run.font.bold = True
        
        # Add project name
        project_box = slide.shapes.add_textbox(
            self.margin, Inches(4.5),
            self.content_width, Inches(1)
        )
        project_frame = project_box.text_frame
        project_frame.text = project_details['project_name']
        project_para = project_frame.paragraphs[0]
        project_para.alignment = PP_ALIGN.CENTER
        project_run = project_para.runs[0]
        project_run.font.size = Pt(32)
        project_run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add client and date
        details_box = slide.shapes.add_textbox(
            self.margin, Inches(6),
            self.content_width, Inches(1)
        )
        details_frame = details_box.text_frame
        details_frame.text = f"{project_details['client_name']} • {project_details['date_range']}"
        details_para = details_frame.paragraphs[0]
        details_para.alignment = PP_ALIGN.CENTER
        details_run = details_para.runs[0]
        details_run.font.size = Pt(20)
        details_run.font.color.rgb = RGBColor(200, 200, 200)
    
    def _create_summary_slide(self, prs: Presentation, project_details: Dict[str, Any], image_analyses: List[Dict[str, Any]]):
        """Create the weekly summary slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin, self.margin,
            self.content_width, Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Weekly Summary"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.size = Pt(36)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(44, 62, 80)
        
        # Add highlights if provided
        if project_details.get('highlights'):
            highlights_box = slide.shapes.add_textbox(
                self.margin, Inches(1.5),
                self.content_width, Inches(1.5)
            )
            highlights_frame = highlights_box.text_frame
            highlights_frame.text = "Key Highlights:"
            highlights_para = highlights_frame.paragraphs[0]
            highlights_para.font.bold = True
            highlights_para.font.size = Pt(20)
            highlights_para.font.color.rgb = RGBColor(44, 62, 80)
            
            # Add highlights text
            highlights_frame.add_paragraph()
            highlights_frame.paragraphs[1].text = project_details['highlights']
            highlights_frame.paragraphs[1].font.size = Pt(16)
            highlights_frame.paragraphs[1].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add progress overview
        progress_box = slide.shapes.add_textbox(
            self.margin, Inches(3.5),
            self.content_width, Inches(3)
        )
        progress_frame = progress_box.text_frame
        progress_frame.text = "Progress Overview:"
        progress_para = progress_frame.paragraphs[0]
        progress_para.font.bold = True
        progress_para.font.size = Pt(20)
        progress_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Calculate overall progress
        total_progress = sum(analysis.get('completion_percentage', 0) for analysis in image_analyses)
        avg_progress = total_progress / len(image_analyses) if image_analyses else 0
        
        # Add progress summary
        progress_frame.add_paragraph()
        progress_frame.paragraphs[1].text = f"• Total Images Analyzed: {len(image_analyses)}"
        progress_frame.paragraphs[1].font.size = Pt(16)
        progress_frame.paragraphs[1].font.color.rgb = RGBColor(52, 73, 94)
        
        progress_frame.add_paragraph()
        progress_frame.paragraphs[2].text = f"• Average Completion: {avg_progress:.1f}%"
        progress_frame.paragraphs[2].font.size = Pt(16)
        progress_frame.paragraphs[2].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add work types summary
        work_types = set(analysis.get('work_type', 'General') for analysis in image_analyses)
        if work_types:
            progress_frame.add_paragraph()
            progress_frame.paragraphs[3].text = f"• Work Types: {', '.join(work_types)}"
            progress_frame.paragraphs[3].font.size = Pt(16)
            progress_frame.paragraphs[3].font.color.rgb = RGBColor(52, 73, 94)
    
    def _create_content_slide(self, prs: Presentation, analysis: Dict[str, Any]):
        """Create a content slide for an analyzed image"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add slide title
        title_box = slide.shapes.add_textbox(
            self.margin, self.margin,
            self.content_width, Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = analysis.get('slide_title', 'Progress Update')
        title_para = title_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(28)
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Add image placeholder (since we can't embed actual images in this demo)
        # In production, you would insert the actual image here
        image_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.margin, Inches(1.5),
            Inches(6), Inches(4)
        )
        image_placeholder.fill.solid()
        image_placeholder.fill.fore_color.rgb = RGBColor(236, 240, 241)  # Light gray
        image_placeholder.line.color.rgb = RGBColor(189, 195, 199)
        image_placeholder.line.width = Pt(2)
        
        # Add image label
        image_label = slide.shapes.add_textbox(
            self.margin, Inches(5.6),
            Inches(6), Inches(0.4)
        )
        image_label_frame = image_label.text_frame
        image_label_frame.text = f"Image: {analysis.get('filename', 'Unknown')}"
        image_label_para = image_label_frame.paragraphs[0]
        image_label_para.alignment = PP_ALIGN.CENTER
        image_label_para.font.size = Pt(12)
        image_label_para.font.color.rgb = RGBColor(127, 140, 141)
        
        # Add content on the right side
        content_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.5),
            Inches(5.5), Inches(4.5)
        )
        content_frame = content_box.text_frame
        
        # Add status
        content_frame.text = f"Status: {analysis.get('progress_status', 'In Progress')}"
        status_para = content_frame.paragraphs[0]
        status_para.font.bold = True
        status_para.font.size = Pt(16)
        status_para.font.color.rgb = RGBColor(46, 204, 113)  # Green
        
        # Add work type
        content_frame.add_paragraph()
        content_frame.paragraphs[1].text = f"Work Type: {analysis.get('work_type', 'General Construction')}"
        content_frame.paragraphs[1].font.size = Pt(14)
        content_frame.paragraphs[1].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add completion percentage
        content_frame.add_paragraph()
        completion = analysis.get('completion_percentage', 50)
        content_frame.paragraphs[2].text = f"Completion: {completion}%"
        content_frame.paragraphs[2].font.size = Pt(14)
        content_frame.paragraphs[2].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add description
        content_frame.add_paragraph()
        content_frame.paragraphs[3].text = "Description:"
        content_frame.paragraphs[3].font.bold = True
        content_frame.paragraphs[3].font.size = Pt(14)
        content_frame.paragraphs[3].font.color.rgb = RGBColor(44, 62, 80)
        
        content_frame.add_paragraph()
        content_frame.paragraphs[4].text = analysis.get('description', 'No description available.')
        content_frame.paragraphs[4].font.size = Pt(12)
        content_frame.paragraphs[4].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add notes if available
        notes = analysis.get('notes', '')
        if notes and notes.strip():
            content_frame.add_paragraph()
            content_frame.paragraphs[5].text = "Notes:"
            content_frame.paragraphs[5].font.bold = True
            content_frame.paragraphs[5].font.size = Pt(14)
            content_frame.paragraphs[5].font.color.rgb = RGBColor(231, 76, 60)  # Red for notes
            
            content_frame.add_paragraph()
            content_frame.paragraphs[6].text = notes
            content_frame.paragraphs[6].font.size = Pt(12)
            content_frame.paragraphs[6].font.color.rgb = RGBColor(52, 73, 94)
    
    def _create_conclusion_slide(self, prs: Presentation, project_details: Dict[str, Any]):
        """Create the conclusion slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            self.margin, Inches(2),
            self.content_width, Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Next Steps & Looking Ahead"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.bold = True
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        
        # Add content
        content_box = slide.shapes.add_textbox(
            self.margin, Inches(3.5),
            self.content_width, Inches(2.5)
        )
        content_frame = content_box.text_frame
        
        content_frame.text = "• Planning for next week's activities is underway"
        content_para = content_frame.paragraphs[0]
        content_para.alignment = PP_ALIGN.CENTER
        content_para.font.size = Pt(20)
        content_para.font.color.rgb = RGBColor(52, 73, 94)
        
        content_frame.add_paragraph()
        content_frame.paragraphs[1].text = "• Regular progress updates will continue"
        content_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        content_frame.paragraphs[1].font.size = Pt(20)
        content_frame.paragraphs[1].font.color.rgb = RGBColor(52, 73, 94)
        
        content_frame.add_paragraph()
        content_frame.paragraphs[2].text = "• Please let us know if you have any questions"
        content_frame.paragraphs[2].alignment = PP_ALIGN.CENTER
        content_frame.paragraphs[2].font.size = Pt(20)
        content_frame.paragraphs[2].font.color.rgb = RGBColor(52, 73, 94)
        
        # Add footer
        footer_box = slide.shapes.add_textbox(
            self.margin, Inches(6.5),
            self.content_width, Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(127, 140, 141)
