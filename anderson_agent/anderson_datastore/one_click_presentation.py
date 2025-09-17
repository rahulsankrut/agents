#!/usr/bin/env python3
"""
One-Click Combined Presentation Generator
Automatically generates a single PowerPoint presentation with all projects
"""

import sys
import os
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from firestore_operations import FirestoreManager
import requests
import json
from datetime import datetime
from google.cloud import storage
from pptx import Presentation
from pptx.util import Inches
import tempfile

# Configuration
CLOUD_FUNCTION_URL = "https://us-central1-agent-space-465923.cloudfunctions.net/ppt-generator"
OUTPUT_BUCKET = "agent-space-465923-presentations"
PROJECT_ID = "agent-space-465923"

def convert_firestore_to_cloud_function_format(project):
    """Convert Firestore project to cloud function format"""
    
    image_data = []
    for img in project.images:
        image_data.append({
            "gcs_url": img.image_url,
            "title": img.description
        })
    
    include_eqi = project.eqi == "Yes"
    
    logo_gcs_url = project.customer_logo_url
    if logo_gcs_url.startswith("https://storage.cloud.google.com/"):
        logo_gcs_url = logo_gcs_url.replace("https://storage.cloud.google.com/", "gs://")
    
    return {
        "title": project.project_title,
        "customer_name": project.customer_name,
        "logo_gcs_url": logo_gcs_url,
        "text_content": [project.project_overview],
        "image_data": image_data,
        "include_eqi": include_eqi
    }

def generate_combined_presentation():
    """Generate a single combined presentation with all projects - ONE CLICK SOLUTION"""
    
    print('🎯 ONE-CLICK COMBINED PRESENTATION GENERATOR')
    print('=' * 60)
    
    # Initialize Firestore manager
    db_manager = FirestoreManager(project_id=PROJECT_ID)
    all_projects = db_manager.list_projects()
    
    print(f'📊 Found {len(all_projects)} projects in database')
    print(f'🔄 Generating combined presentation...')
    
    if len(all_projects) == 0:
        print('❌ No projects found in database')
        return None
    
    # Create a new presentation
    combined_prs = Presentation()
    combined_prs.slide_width = Inches(13.33)
    combined_prs.slide_height = Inches(7.5)
    
    slide_count = 0
    failed_count = 0
    
    # Process each project
    for i, project in enumerate(all_projects, 1):
        try:
            print(f'📝 Processing project {i}/{len(all_projects)}: {project.project_title[:40]}...')
            
            # Convert to cloud function format
            project_data = convert_firestore_to_cloud_function_format(project)
            
            # Generate individual slide
            response = requests.post(
                CLOUD_FUNCTION_URL + "/generate",
                json=project_data,
                headers={'Content-Type': 'application/json'},
                timeout=300
            )
            
            if response.status_code == 200:
                # Save individual slide to temporary file
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                    temp_file.write(response.content)
                    temp_file_path = temp_file.name
                
                # Open the individual slide
                individual_prs = Presentation(temp_file_path)
                
                # Copy each slide from the individual presentation
                for slide in individual_prs.slides:
                    slide_layout = combined_prs.slide_layouts[6]  # Blank layout
                    new_slide = combined_prs.slides.add_slide(slide_layout)
                    
                    # Copy all shapes from the original slide
                    for shape in slide.shapes:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Copy the shape based on its type
                        if hasattr(shape, 'image'):
                            try:
                                new_slide.shapes.add_picture(
                                    shape.image.blob, 
                                    left, top, width, height
                                )
                            except:
                                pass
                        elif hasattr(shape, 'text_frame'):
                            try:
                                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                                text_frame = textbox.text_frame
                                
                                if shape.text_frame.paragraphs:
                                    for paragraph in shape.text_frame.paragraphs:
                                        p = text_frame.add_paragraph()
                                        p.text = paragraph.text
                                        p.level = paragraph.level
                                        
                                        if paragraph.font:
                                            p.font.name = paragraph.font.name
                                            p.font.size = paragraph.font.size
                                            p.font.bold = paragraph.font.bold
                                            p.font.italic = paragraph.font.italic
                            except:
                                pass
                    
                    slide_count += 1
                
                # Clean up temporary file
                os.unlink(temp_file_path)
                
            else:
                print(f'❌ Failed to generate slide for project {i}')
                failed_count += 1
                
        except Exception as e:
            print(f'❌ Error processing project {i}: {e}')
            failed_count += 1
            continue
    
    # Save combined presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Combined_Presentation_{timestamp}.pptx"
    
    try:
        # Save locally
        combined_prs.save(filename)
        print(f'✅ Combined presentation saved locally: {filename}')
        
        # Save to GCS
        client = storage.Client(project=PROJECT_ID)
        bucket = client.bucket(OUTPUT_BUCKET)
        blob = bucket.blob(filename)
        blob.upload_from_filename(filename)
        
        print(f'✅ Combined presentation saved to GCS: gs://{OUTPUT_BUCKET}/{filename}')
        
        print(f'\n🎉 SUCCESS!')
        print(f'📁 File: {filename}')
        print(f'📊 Total slides: {slide_count}')
        print(f'✅ Successful: {slide_count}')
        print(f'❌ Failed: {failed_count}')
        print(f'📂 Saved locally and to Cloud Storage')
        
        return filename
        
    except Exception as e:
        print(f'❌ Error saving combined presentation: {e}')
        return None

if __name__ == "__main__":
    generate_combined_presentation()
