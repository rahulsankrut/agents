#!/usr/bin/env python3
"""
Automated PowerPoint Slide Combiner
Combines individual slides into a single presentation seamlessly
"""

import sys
import os
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from firestore_operations import FirestoreManager
import requests
import json
from datetime import datetime
from google.cloud import storage
from pptx import Presentation
from pptx.util import Inches
import tempfile

# Configuration
CLOUD_FUNCTION_URL = "https://us-central1-agent-space-465923.cloudfunctions.net/ppt-generator"
OUTPUT_BUCKET = "agent-space-465923-presentations"
PROJECT_ID = "agent-space-465923"

def convert_firestore_to_cloud_function_format(project):
    """Convert Firestore project to cloud function format"""
    
    # Convert image data
    image_data = []
    for img in project.images:
        image_data.append({
            "gcs_url": img.image_url,
            "title": img.description
        })
    
    # Convert EQI from Yes/No to boolean
    include_eqi = project.eqi == "Yes"
    
    # Convert customer logo URL to GCS format if needed
    logo_gcs_url = project.customer_logo_url
    if logo_gcs_url.startswith("https://storage.cloud.google.com/"):
        logo_gcs_url = logo_gcs_url.replace("https://storage.cloud.google.com/", "gs://")
    
    return {
        "title": project.project_title,
        "logo_gcs_url": logo_gcs_url,
        "text_content": [project.project_overview],
        "image_data": image_data,
        "include_eqi": include_eqi
    }

def generate_individual_slide(project_data):
    """Generate a single slide using the cloud function"""
    
    try:
        print(f"🔄 Generating slide: {project_data['title'][:50]}...")
        
        # Make the request
        response = requests.post(
            CLOUD_FUNCTION_URL + "/generate",
            json=project_data,
            headers={'Content-Type': 'application/json'},
            timeout=300
        )
        
        if response.status_code == 200:
            pptx_data = response.content
            print(f"✅ Generated slide ({len(pptx_data)} bytes)")
            return pptx_data
        else:
            print(f"❌ Error: {response.status_code} - {response.text}")
            return None
            
    except Exception as e:
        print(f"❌ Error generating slide: {e}")
        return None

def combine_slides_into_single_presentation(slide_data_list):
    """Combine multiple PowerPoint slides into a single presentation"""
    
    print('🔗 COMBINING SLIDES INTO SINGLE PRESENTATION')
    print('=' * 50)
    
    # Create a new presentation
    combined_prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    combined_prs.slide_width = Inches(13.33)
    combined_prs.slide_height = Inches(7.5)
    
    slide_count = 0
    
    for i, slide_data in enumerate(slide_data_list, 1):
        if slide_data is None:
            print(f"⚠️  Skipping slide {i} (no data)")
            continue
            
        try:
            # Save individual slide to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(slide_data)
                temp_file_path = temp_file.name
            
            # Open the individual slide
            individual_prs = Presentation(temp_file_path)
            
            # Copy each slide from the individual presentation
            for slide in individual_prs.slides:
                # Get the slide layout (use blank layout)
                slide_layout = combined_prs.slide_layouts[6]  # Blank layout
                new_slide = combined_prs.slides.add_slide(slide_layout)
                
                # Copy all shapes from the original slide
                for shape in slide.shapes:
                    # Get the shape's position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Copy the shape based on its type
                    if hasattr(shape, 'image'):
                        # It's a picture
                        try:
                            new_slide.shapes.add_picture(
                                shape.image.blob, 
                                left, top, width, height
                            )
                        except:
                            # If we can't copy the image, skip it
                            pass
                    elif hasattr(shape, 'text_frame'):
                        # It's a text box
                        try:
                            textbox = new_slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            
                            # Copy text content
                            if shape.text_frame.paragraphs:
                                for paragraph in shape.text_frame.paragraphs:
                                    p = text_frame.add_paragraph()
                                    p.text = paragraph.text
                                    p.level = paragraph.level
                                    
                                    # Copy formatting
                                    if paragraph.font:
                                        p.font.name = paragraph.font.name
                                        p.font.size = paragraph.font.size
                                        p.font.bold = paragraph.font.bold
                                        p.font.italic = paragraph.font.italic
                        except:
                            # If we can't copy the text, skip it
                            pass
                
                slide_count += 1
                print(f"✅ Added slide {slide_count}")
            
            # Clean up temporary file
            os.unlink(temp_file_path)
            
        except Exception as e:
            print(f"❌ Error processing slide {i}: {e}")
            continue
    
    print(f'🎉 Successfully combined {slide_count} slides into single presentation')
    return combined_prs

def save_combined_presentation(combined_prs, filename):
    """Save the combined presentation locally and to GCS"""
    
    try:
        # Save locally
        combined_prs.save(filename)
        print(f"✅ Combined presentation saved locally: {filename}")
        
        # Save to GCS
        client = storage.Client(project=PROJECT_ID)
        bucket = client.bucket(OUTPUT_BUCKET)
        
        blob = bucket.blob(filename)
        blob.upload_from_filename(filename)
        
        print(f"✅ Combined presentation saved to GCS: gs://{OUTPUT_BUCKET}/{filename}")
        
        return filename
        
    except Exception as e:
        print(f"❌ Error saving combined presentation: {e}")
        return None

def generate_combined_presentation():
    """Generate a single combined presentation with all projects"""
    
    print('🎯 GENERATING COMBINED PRESENTATION')
    print('=' * 50)
    
    # Initialize Firestore manager
    db_manager = FirestoreManager(project_id=PROJECT_ID)
    
    # Get all projects
    all_projects = db_manager.list_projects()
    print(f'📊 Found {len(all_projects)} projects in database')
    
    if len(all_projects) == 0:
        print('❌ No projects found in database')
        return
    
    # Generate individual slides
    slide_data_list = []
    failed_projects = []
    
    print(f'\n🔄 Generating {len(all_projects)} individual slides...')
    
    for i, project in enumerate(all_projects, 1):
        print(f'\n📝 Processing project {i}/{len(all_projects)}')
        print(f'   Title: {project.project_title[:60]}...')
        print(f'   Customer: {project.customer_name}')
        print(f'   EQI: {project.eqi}')
        print(f'   Images: {len(project.images)}')
        
        # Convert to cloud function format
        project_data = convert_firestore_to_cloud_function_format(project)
        
        # Generate individual slide
        slide_data = generate_individual_slide(project_data)
        
        if slide_data:
            slide_data_list.append(slide_data)
        else:
            slide_data_list.append(None)
            failed_projects.append(project.project_title)
    
    print(f'\n📊 Generated {len([s for s in slide_data_list if s is not None])} slides')
    if failed_projects:
        print(f'⚠️  Failed projects: {len(failed_projects)}')
        for failed in failed_projects:
            print(f'   - {failed}')
    
    # Combine slides into single presentation
    if slide_data_list:
        combined_prs = combine_slides_into_single_presentation(slide_data_list)
        
        # Save combined presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"Combined_Presentation_{timestamp}.pptx"
        
        saved_file = save_combined_presentation(combined_prs, filename)
        
        if saved_file:
            print(f'\n🎉 COMBINED PRESENTATION GENERATED SUCCESSFULLY!')
            print(f'📁 File: {filename}')
            print(f'📊 Total slides: {len(combined_prs.slides)}')
            print(f'📂 Saved locally and to GCS')
            
            # Create summary
            summary_data = {
                "metadata": {
                    "generated_at": datetime.now().isoformat(),
                    "total_projects": len(all_projects),
                    "successful_slides": len([s for s in slide_data_list if s is not None]),
                    "failed_projects": failed_projects,
                    "combined_presentation": filename,
                    "gcs_location": f"gs://{OUTPUT_BUCKET}/{filename}"
                }
            }
            
            summary_filename = f"combined_presentation_summary_{timestamp}.json"
            with open(summary_filename, 'w') as f:
                json.dump(summary_data, f, indent=2)
            
            print(f'📋 Summary saved to: {summary_filename}')
        else:
            print('❌ Failed to save combined presentation')
    else:
        print('❌ No slides generated to combine')

def main():
    """Main function"""
    
    print('🎯 AUTOMATED PRESENTATION COMBINER')
    print('=' * 40)
    print(f'Cloud Function: {CLOUD_FUNCTION_URL}')
    print(f'Output Bucket: gs://{OUTPUT_BUCKET}')
    print(f'Project ID: {PROJECT_ID}')
    print()
    
    print('This will:')
    print('1. Generate individual slides for all projects')
    print('2. Combine them into a single PowerPoint presentation')
    print('3. Save locally and to Cloud Storage')
    print()
    
    choice = input('Generate combined presentation? (y/n): ')
    
    if choice.lower() in ['y', 'yes']:
        generate_combined_presentation()
    else:
        print('👋 Operation cancelled')

if __name__ == "__main__":
    main()
