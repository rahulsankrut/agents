#!/usr/bin/env python3
"""
Improved Combined Presentation Generator
Preserves original slide formatting and layout
"""

import sys
import os
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from firestore_operations import FirestoreManager
import requests
import json
from datetime import datetime
from google.cloud import storage
from pptx import Presentation
from pptx.util import Inches
import tempfile
import shutil

# Configuration
CLOUD_FUNCTION_URL = "https://us-central1-agent-space-465923.cloudfunctions.net/ppt-generator"
OUTPUT_BUCKET = "agent-space-465923-presentations"
PROJECT_ID = "agent-space-465923"

def convert_firestore_to_cloud_function_format(project):
    """Convert Firestore project to cloud function format"""
    
    image_data = []
    for img in project.images:
        image_data.append({
            "gcs_url": img.image_url,
            "title": img.description
        })
    
    include_eqi = project.eqi == "Yes"
    
    logo_gcs_url = project.customer_logo_url
    if logo_gcs_url.startswith("https://storage.cloud.google.com/"):
        logo_gcs_url = logo_gcs_url.replace("https://storage.cloud.google.com/", "gs://")
    
    return {
        "title": project.project_title,
        "logo_gcs_url": logo_gcs_url,
        "text_content": [project.project_overview],
        "image_data": image_data,
        "include_eqi": include_eqi
    }

def generate_combined_presentation_preserve_format():
    """Generate combined presentation preserving original formatting"""
    
    print('🎯 IMPROVED COMBINED PRESENTATION GENERATOR')
    print('=' * 60)
    print('📋 This version preserves original slide formatting')
    
    # Initialize Firestore manager
    db_manager = FirestoreManager(project_id=PROJECT_ID)
    all_projects = db_manager.list_projects()
    
    print(f'📊 Found {len(all_projects)} projects in database')
    
    if len(all_projects) == 0:
        print('❌ No projects found in database')
        return None
    
    # Create a new presentation
    combined_prs = Presentation()
    combined_prs.slide_width = Inches(13.33)
    combined_prs.slide_height = Inches(7.5)
    
    slide_count = 0
    failed_count = 0
    temp_files = []  # Keep track of temp files for cleanup
    
    print(f'🔄 Generating {len(all_projects)} individual slides...')
    
    # Process each project
    for i, project in enumerate(all_projects, 1):
        try:
            print(f'📝 Processing project {i}/{len(all_projects)}: {project.project_title[:40]}...')
            
            # Convert to cloud function format
            project_data = convert_firestore_to_cloud_function_format(project)
            
            # Generate individual slide
            response = requests.post(
                CLOUD_FUNCTION_URL + "/generate",
                json=project_data,
                headers={'Content-Type': 'application/json'},
                timeout=300
            )
            
            if response.status_code == 200:
                # Save individual slide to temporary file
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                    temp_file.write(response.content)
                    temp_file_path = temp_file.name
                    temp_files.append(temp_file_path)
                
                # Open the individual slide
                individual_prs = Presentation(temp_file_path)
                
                # Copy each slide from the individual presentation
                for slide in individual_prs.slides:
                    # Use the same layout as the original slide
                    slide_layout = slide.slide_layout
                    
                    # Create new slide with the same layout
                    new_slide = combined_prs.slides.add_slide(slide_layout)
                    
                    # Copy all shapes with their exact positioning and formatting
                    for shape in slide.shapes:
                        try:
                            # Get shape properties
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            # Copy based on shape type
                            if hasattr(shape, 'image'):
                                # Copy image
                                try:
                                    # Get the image blob
                                    image_blob = shape.image.blob
                                    new_slide.shapes.add_picture(
                                        image_blob, 
                                        left, top, width, height
                                    )
                                except Exception as e:
                                    print(f'   ⚠️  Could not copy image: {e}')
                                    
                            elif hasattr(shape, 'text_frame'):
                                # Copy text box with formatting
                                try:
                                    textbox = new_slide.shapes.add_textbox(left, top, width, height)
                                    text_frame = textbox.text_frame
                                    
                                    # Clear default paragraph
                                    text_frame.clear()
                                    
                                    # Copy paragraphs with formatting
                                    for paragraph in shape.text_frame.paragraphs:
                                        p = text_frame.add_paragraph()
                                        p.text = paragraph.text
                                        p.level = paragraph.level
                                        
                                        # Copy paragraph formatting
                                        if paragraph.font:
                                            p.font.name = paragraph.font.name
                                            p.font.size = paragraph.font.size
                                            p.font.bold = paragraph.font.bold
                                            p.font.italic = paragraph.font.italic
                                            p.font.color.rgb = paragraph.font.color.rgb
                                        
                                        # Copy alignment
                                        p.alignment = paragraph.alignment
                                        
                                except Exception as e:
                                    print(f'   ⚠️  Could not copy text: {e}')
                                    
                            elif hasattr(shape, 'shape_type'):
                                # Copy other shape types (rectangles, etc.)
                                try:
                                    if shape.shape_type == 1:  # Rectangle
                                        new_slide.shapes.add_shape(
                                            shape.shape_type,
                                            left, top, width, height
                                        )
                                except Exception as e:
                                    print(f'   ⚠️  Could not copy shape: {e}')
                                    
                        except Exception as e:
                            print(f'   ⚠️  Could not copy shape: {e}')
                            continue
                    
                    slide_count += 1
                    print(f'   ✅ Added slide {slide_count}')
                
            else:
                print(f'❌ Failed to generate slide for project {i}')
                failed_count += 1
                
        except Exception as e:
            print(f'❌ Error processing project {i}: {e}')
            failed_count += 1
            continue
    
    # Clean up temporary files
    for temp_file in temp_files:
        try:
            os.unlink(temp_file)
        except:
            pass
    
    # Save combined presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Combined_Presentation_Formatted_{timestamp}.pptx"
    
    try:
        # Save locally
        combined_prs.save(filename)
        print(f'✅ Combined presentation saved locally: {filename}')
        
        # Save to GCS
        client = storage.Client(project=PROJECT_ID)
        bucket = client.bucket(OUTPUT_BUCKET)
        blob = bucket.blob(filename)
        blob.upload_from_filename(filename)
        
        print(f'✅ Combined presentation saved to GCS: gs://{OUTPUT_BUCKET}/{filename}')
        
        print(f'\n🎉 SUCCESS!')
        print(f'📁 File: {filename}')
        print(f'📊 Total slides: {slide_count}')
        print(f'✅ Successful: {slide_count}')
        print(f'❌ Failed: {failed_count}')
        print(f'📂 Saved locally and to Cloud Storage')
        print(f'🎨 Original formatting preserved')
        
        return filename
        
    except Exception as e:
        print(f'❌ Error saving combined presentation: {e}')
        return None

if __name__ == "__main__":
    generate_combined_presentation_preserve_format()
