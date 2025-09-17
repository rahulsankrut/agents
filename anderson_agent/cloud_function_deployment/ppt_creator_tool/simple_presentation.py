#!/usr/bin/env python3
"""
Simple PowerPoint Presentation Generator
Creates a basic PowerPoint presentation with a header image and content.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_simple_presentation(title, logo_path=None, text_content=None, image_data=None, include_eqi=False):
    """Create a simple PowerPoint presentation with header image"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slide with blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add header image if it exists
    header_image_path = "slides_stateful_resources/slide_header.png"
    if os.path.exists(header_image_path):
        try:
            # Add header image at the top
            slide.shapes.add_picture(header_image_path, Inches(0), Inches(0), Inches(13.33), Inches(1.5))
            print(f"✅ Added header image: {header_image_path}")
        except Exception as e:
            print(f"❌ Could not add header image: {e}")
            create_text_header(slide)
    else:
        print(f"⚠️  Header image not found: {header_image_path}")
        create_text_header(slide)
    
    # Add logo outside and to the left of the title box if provided
    if logo_path and os.path.exists(logo_path):
        try:
            # Get original logo dimensions to maintain aspect ratio
            from PIL import Image
            with Image.open(logo_path) as img:
                original_width, original_height = img.size
                aspect_ratio = original_width / original_height
            
            # Define maximum dimensions for the logo
            max_width = Inches(0.8)
            max_height = Inches(0.8)
            
            # Calculate dimensions maintaining aspect ratio
            if max_width / max_height > aspect_ratio:
                # Height is the limiting factor
                new_height = max_height
                new_width = new_height * aspect_ratio
            else:
                # Width is the limiting factor
                new_width = max_width
                new_height = new_width / aspect_ratio
            
            # Add logo with maintained aspect ratio
            slide.shapes.add_picture(logo_path, Inches(0.2), Inches(2.1), new_width, new_height)
            print(f"✅ Added logo: {logo_path} (maintained aspect ratio: {new_width:.2f}\" x {new_height:.2f}\")")
            
            # Add slide sub-header image beside the logo ONLY if EQI is included
            if include_eqi:
                sub_header_image_path = "slides_stateful_resources/slide_sub_header.png"
                if os.path.exists(sub_header_image_path):
                    try:
                        # Get original image dimensions to maintain aspect ratio
                        from PIL import Image
                        with Image.open(sub_header_image_path) as img:
                            original_width, original_height = img.size
                            aspect_ratio = original_width / original_height
                        
                        # Define maximum dimensions for the sub-header
                        max_width = Inches(2.0)
                        max_height = Inches(0.8)
                        
                        # Calculate dimensions maintaining aspect ratio
                        if max_width / max_height > aspect_ratio:
                            # Height is the limiting factor
                            new_height = max_height
                            new_width = new_height * aspect_ratio
                        else:
                            # Width is the limiting factor
                            new_width = max_width
                            new_height = new_width / aspect_ratio
                        
                        sub_header_x = Inches(1.1)  # Position beside logo
                        sub_header_y = Inches(2.1)
                        
                        # Add sub-header image with maintained aspect ratio
                        slide.shapes.add_picture(sub_header_image_path, sub_header_x, sub_header_y, new_width, new_height)
                        print(f"✅ Added sub-header image: {sub_header_image_path} (maintained aspect ratio: {new_width:.2f}\" x {new_height:.2f}\")")
                        
                        # Adjust title box position to make room for logo and sub-header (using actual sub-header width)
                        title_rect_x = sub_header_x + new_width + Inches(0.1)  # Start after logo + sub-header + small gap
                        title_rect_width = Inches(13.33) - title_rect_x - Inches(0.2)  # Remaining width minus margins
                    except Exception as e:
                        print(f"❌ Could not add sub-header image: {e}")
                        # Fall back to logo-only layout
                        title_rect_x = Inches(1.1)
                        title_rect_width = Inches(11.23)
                else:
                    print(f"⚠️  Sub-header image not found: {sub_header_image_path}")
                    # Fall back to logo-only layout
                    title_rect_x = Inches(1.1)
                    title_rect_width = Inches(11.23)
            else:
                # No sub-header, adjust title box position to make room for logo only
                title_rect_x = Inches(1.1)  # Start title box after logo
                title_rect_width = Inches(11.23)  # Adjust width to account for logo space only
        except Exception as e:
            print(f"❌ Could not add logo: {e}")
            # Fall back to full width title box
            title_rect_x = Inches(1)
            title_rect_width = Inches(11.33)
    else:
        if logo_path:
            print(f"⚠️  Logo not found: {logo_path}")
        # No logo, use full width for title box
        title_rect_x = Inches(1)
        title_rect_width = Inches(11.33)
    
    # Add title in rounded rectangle box
    from pptx.enum.shapes import MSO_SHAPE
    title_rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, title_rect_x, Inches(2), title_rect_width, Inches(1)
    )
    title_rect.fill.solid()
    title_rect.fill.fore_color.rgb = RGBColor(189, 215, 238)  # Light blue 5 color
    title_rect.line.color.rgb = RGBColor(0, 51, 102)  # Dark blue border
    title_rect.line.width = Pt(2)
    
    # Add title text inside the title box
    title_box = slide.shapes.add_textbox(title_rect_x + Inches(0.1), Inches(2.1), title_rect_width - Inches(0.2), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.name = 'Calibri'
    title_p.font.size = Pt(19.2)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue text
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add content section with two rectangle boxes
    # Left box - Text content (3 inches width)
    left_box_x = Inches(1)
    left_box_y = Inches(3.5)
    left_box_width = Inches(3)
    left_box_height = Inches(3)
    
    # Right box - Images (remaining width)
    right_box_x = Inches(4.5)
    right_box_y = Inches(3.5)
    right_box_width = Inches(7.83)
    right_box_height = Inches(3)
    
    # Create left rectangle box for text content
    from pptx.enum.shapes import MSO_SHAPE
    left_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left_box_x, left_box_y, left_box_width, left_box_height
    )
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    left_rect.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    left_rect.line.width = Pt(3)  # Thick border
    
    # Add heading to left box
    heading_box = slide.shapes.add_textbox(left_box_x + Inches(0.1), left_box_y + Inches(0.1), left_box_width - Inches(0.2), Inches(0.4))
    heading_frame = heading_box.text_frame
    heading_frame.clear()
    heading_p = heading_frame.paragraphs[0]
    heading_p.text = "Project Overview/Callouts"
    heading_p.font.name = 'Calibri'
    heading_p.font.size = Pt(12.8)
    heading_p.font.bold = True
    heading_p.font.color.rgb = RGBColor(0, 0, 0)
    heading_p.alignment = PP_ALIGN.LEFT
    
    # Add text content to left box
    text_content_box = slide.shapes.add_textbox(left_box_x + Inches(0.1), left_box_y + Inches(0.6), left_box_width - Inches(0.2), left_box_height - Inches(0.7))
    text_content_frame = text_content_box.text_frame
    text_content_frame.clear()
    text_content_frame.word_wrap = True
    
    if text_content:
        # Handle both string and list inputs
        if isinstance(text_content, list):
            lines = text_content
        else:
            # Split content by lines and add as bullet points
            lines = text_content.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = text_content_frame.paragraphs[0]
            else:
                p = text_content_frame.add_paragraph()
            
            line = line.strip()
            if line:
                # Add bullet point if not already present
                if not line.startswith('•'):
                    line = '• ' + line
                p.text = line
                p.font.name = 'Calibri'
                p.font.size = Pt(11.2)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT
    else:
        # Default content
        default_p = text_content_frame.paragraphs[0]
        default_p.text = "• Add your project content here\n• Use bullet points for clarity\n• Include key information"
        default_p.font.name = 'Calibri'
        default_p.font.size = Pt(11.2)
        default_p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Create right rectangle box for images
    right_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, right_box_x, right_box_y, right_box_width, right_box_height
    )
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    right_rect.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    right_rect.line.width = Pt(3)  # Thick border
    
    # Add images to right box
    if image_data and isinstance(image_data, list):
        num_images = len(image_data)
        if num_images > 0:
            # Calculate available space for images
            available_width = right_box_width - Inches(0.4)  # Account for padding
            available_height = right_box_height - Inches(0.4)  # Account for padding
            space_per_image = available_width / num_images
            
            for i, img_data in enumerate(image_data):
                # Extract image path and title
                if isinstance(img_data, dict):
                    image_path = img_data.get('path', '')
                    image_title = img_data.get('title', f'Image {i+1}')
                else:
                    # Fallback for old format (just path)
                    image_path = img_data
                    image_title = f'Image {i+1}'
                
                if image_path and os.path.exists(image_path):
                    try:
                        # Get original image dimensions to maintain aspect ratio
                        from PIL import Image
                        with Image.open(image_path) as img:
                            original_width, original_height = img.size
                            aspect_ratio = original_width / original_height
                        
                        # Calculate dimensions maintaining aspect ratio
                        max_width = space_per_image - Inches(0.1)  # Leave some space between images
                        max_height = available_height - Inches(0.3)  # Leave space for title
                        
                        # Calculate new dimensions maintaining aspect ratio
                        if max_width / max_height > aspect_ratio:
                            # Height is the limiting factor
                            new_height = max_height
                            new_width = new_height * aspect_ratio
                        else:
                            # Width is the limiting factor
                            new_width = max_width
                            new_height = new_width / aspect_ratio
                        
                        # Position images side by side
                        image_x = right_box_x + Inches(0.2) + (i * space_per_image)
                        image_y = right_box_y + Inches(0.5)  # Leave space for title
                        
                        # Add image title
                        title_box = slide.shapes.add_textbox(image_x, right_box_y + Inches(0.2), space_per_image - Inches(0.1), Inches(0.25))
                        title_frame = title_box.text_frame
                        title_frame.clear()
                        title_p = title_frame.paragraphs[0]
                        title_p.text = image_title
                        title_p.font.name = 'Calibri'
                        title_p.font.size = Pt(9.6)
                        title_p.font.bold = True
                        title_p.font.color.rgb = RGBColor(0, 0, 0)
                        title_p.alignment = PP_ALIGN.CENTER
                        
                        # Add image
                        slide.shapes.add_picture(image_path, image_x, image_y, new_width, new_height)
                        print(f"✅ Added image: {image_path} with title: {image_title} (maintained aspect ratio)")
                        print(f"   Image file exists: {os.path.exists(image_path)}")
                        print(f"   Image file size: {os.path.getsize(image_path)} bytes")
                        print(f"   Image dimensions: {original_width}x{original_height}")
                        print(f"   Final dimensions: {new_width}x{new_height}")
                    except Exception as e:
                        print(f"❌ Could not add image {image_path}: {e}")
                        # Add placeholder with default dimensions
                        placeholder_width = space_per_image - Inches(0.1)
                        placeholder_height = available_height - Inches(0.3)
                        image_x = right_box_x + Inches(0.2) + (i * space_per_image)
                        image_y = right_box_y + Inches(0.5)
                        add_image_placeholder(slide, image_x, image_y, placeholder_width, placeholder_height, f"Image {i+1}")
                else:
                    print(f"⚠️  Image not found: {image_path}")
                    # Add placeholder with default dimensions
                    placeholder_width = space_per_image - Inches(0.1)
                    placeholder_height = available_height - Inches(0.3)
                    image_x = right_box_x + Inches(0.2) + (i * space_per_image)
                    image_y = right_box_y + Inches(0.5)
                    add_image_placeholder(slide, image_x, image_y, placeholder_width, placeholder_height, f"Image {i+1}")
    else:
        # Add placeholder text
        placeholder_box = slide.shapes.add_textbox(right_box_x + Inches(0.2), right_box_y + right_box_height/2 - Inches(0.1), right_box_width - Inches(0.4), Inches(0.2))
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.clear()
        placeholder_p = placeholder_frame.paragraphs[0]
        placeholder_p.text = "Images"
        placeholder_p.font.name = 'Calibri'
        placeholder_p.font.size = Pt(11.2)
        placeholder_p.font.color.rgb = RGBColor(128, 128, 128)
        placeholder_p.alignment = PP_ALIGN.CENTER
    
    # Add footer with proprietary notice
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.33), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Add your footer text here"
    footer_p.font.name = 'Calibri'
    footer_p.font.size = Pt(8)
    footer_p.font.color.rgb = RGBColor(128, 128, 128)
    footer_p.alignment = PP_ALIGN.CENTER
    
    return prs

def create_text_header(slide):
    """Create a text-based header if image is not available"""
    # Create header background
    from pptx.enum.shapes import MSO_SHAPE
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.5)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
    header_bg.line.fill.background()
    
    # Add header text
    header_text_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(13.33), Inches(0.5))
    header_text_frame = header_text_box.text_frame
    header_text_frame.clear()
    header_text_p = header_text_frame.paragraphs[0]
    header_text_p.text = "Project Header"
    header_text_p.font.name = 'Calibri'
    header_text_p.font.size = Pt(19.2)
    header_text_p.font.bold = True
    header_text_p.font.color.rgb = RGBColor(0, 51, 102)
    header_text_p.alignment = PP_ALIGN.CENTER

def add_image_placeholder(slide, x, y, width, height, label):
    """Add an image placeholder with label"""
    # Add placeholder rectangle
    from pptx.enum.shapes import MSO_SHAPE
    placeholder_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, width, height
    )
    placeholder_rect.fill.solid()
    placeholder_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
    placeholder_rect.line.color.rgb = RGBColor(128, 128, 128)  # Gray border
    placeholder_rect.line.width = Pt(1)
    
    # Add placeholder text
    placeholder_box = slide.shapes.add_textbox(x, y, width, height)
    placeholder_frame = placeholder_box.text_frame
    placeholder_frame.clear()
    placeholder_p = placeholder_frame.paragraphs[0]
    placeholder_p.text = label
    placeholder_p.font.name = 'Calibri'
    placeholder_p.font.size = Pt(9.6)
    placeholder_p.font.color.rgb = RGBColor(128, 128, 128)
    placeholder_p.alignment = PP_ALIGN.CENTER

def main():
    """Main function to create and save presentation"""
    try:
        print("Creating simple PowerPoint presentation...")
        
        # Get title from user input
        print("\nEnter the title for your presentation:")
        title = input("Title: ").strip()
        
        if not title:
            print("❌ Title is required. Please provide a title.")
            return
        
        # Get logo path from user input
        print("\nEnter the path to your logo image (or press Enter to skip):")
        logo_path = input("Logo path: ").strip()
        
        # Get text content from user input
        print("\nEnter text content for the left box (press Enter twice when done):")
        text_lines = []
        while True:
            line = input()
            if line == "" and text_lines and text_lines[-1] == "":
                break
            text_lines.append(line)
        
        text_content = "\n".join(text_lines[:-1]) if text_lines else None
        
        # Get image data from user input
        print("\nEnter image data for the right box (press Enter to finish):")
        image_data = []
        while True:
            image_path = input("Image path: ").strip()
            if not image_path:
                break
            image_title = input("Image title: ").strip()
            if not image_title:
                image_title = f"Image {len(image_data) + 1}"
            image_data.append({"path": image_path, "title": image_title})
        
        # Ask about Execution Quality Index
        print("\nInclude 'Execution Quality Index' in sub-header? (y/n):")
        eqi_input = input("Include EQI: ").strip().lower()
        include_eqi = eqi_input in ['y', 'yes', '1', 'true']
        
        # Create presentation
        prs = create_simple_presentation(title, logo_path, text_content, image_data, include_eqi)
        
        # Save presentation
        output_filename = "simple_presentation.pptx"
        prs.save(output_filename)
        
        print(f"✅ Presentation saved as: {output_filename}")
        print("📁 Location: " + os.path.abspath(output_filename))
        
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")

if __name__ == "__main__":
    main()
